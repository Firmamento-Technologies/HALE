"""
Build PPTX deck 3 slide EN:
1. DOPE Hubs + Firmamento Technologies (team photo + ecosystem)
2. Project HALE (intro + tech specs + 4 use cases)
3. Beyond the Drone (3 pillars + Pentema pilot + numbers)

Stile coerente con il deck master Firmamento Technologies (dark navy + gold +
light blue, font sans-serif + monospace).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = "/home/user/HALE/studio-di-fattibilita"
ASSETS = "/home/user/HALE/cad"
OUT = os.path.join(BASE, "deliverables-FINAL")
LOGO = os.path.join(ASSETS, "LogoFirmamento Technologies.png")
HALE_RENDER = os.path.join(ASSETS, "HALE2.png")
TEAM_PHOTO = os.path.join(ASSETS, "team_firmamento_wide.jpg")
OUT_PPTX = os.path.join(OUT, "DECK-HALE-3slide-EN.pptx")

NAVY = RGBColor(0x0a, 0x1a, 0x3d)
GOLD = RGBColor(0xf0, 0xc9, 0x5c)
GOLD_DARK = RGBColor(0xc9, 0xa6, 0x4a)
LIGHT_BLUE = RGBColor(0xa8, 0xc5, 0xeb)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREY = RGBColor(0x3a, 0x3a, 0x3a)
GREY_LIGHT = RGBColor(0xe6, 0xe6, 0xea)

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

FONT_DISPLAY = "Inter"
FONT_BODY = "Inter"
FONT_MONO = "Consolas"


def text_box(slide, x, y, w, h, text, font_size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=FONT_BODY, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font_name
    return tb


def rich_text_box(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r_def in runs:
        r = p.add_run()
        r.text = r_def["text"]
        r.font.size = Pt(r_def.get("size", 12))
        r.font.bold = r_def.get("bold", False)
        r.font.italic = r_def.get("italic", False)
        r.font.color.rgb = r_def.get("color", WHITE)
        r.font.name = r_def.get("font", FONT_BODY)
    return tb


def add_rect(slide, x, y, w, h, fill=NAVY, line=None, line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GOLD, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(path, x, y, width=w)
    if h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


# ============== SLIDE 1: ECOSYSTEM (DOPE + FIRMAMENTO) ==============
def slide_1_ecosystem():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, fill=NAVY)

    # Team photo: pre-cropped wide format 2400×779 ratio 3.08
    # Slide width SW = 33.867 cm, photo height auto = 33.867 / 3.08 = ~11 cm
    add_image(s, TEAM_PHOTO, 0, 0, w=SW)

    # Tagline overlay top of photo - dark band semi-transparent
    overlay_top = add_rect(s, 0, 0, SW, Cm(3.5), fill=NAVY)
    # transparency via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    sppr = overlay_top.fill._xPr.find(qn("a:solidFill"))
    if sppr is not None:
        clr = sppr.find(qn("a:srgbClr"))
        if clr is not None:
            alpha = etree.SubElement(clr, qn("a:alpha"))
            alpha.set("val", "55000")  # 55% opacity

    text_box(s, Cm(1.5), Cm(0.9), SW - Cm(3), Cm(0.9),
             "DEEP-TECH ECOSYSTEM",
             font_size=13, bold=True, color=GOLD, font_name=FONT_DISPLAY,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    text_box(s, Cm(1.5), Cm(1.6), SW - Cm(3), Cm(1.8),
             "Humans Keep Purpose",
             font_size=40, bold=True, color=WHITE, font_name=FONT_DISPLAY,
             align=PP_ALIGN.CENTER, line_spacing=1.0)

    # No overlap: photo is exactly 11 cm tall now (pre-cropped ratio matches)
    # Add subtle gradient strip under photo
    add_rect(s, 0, Cm(11), SW, Cm(0.1), fill=GOLD)

    # Two columns DOPE + Firmamento
    bottom_y = Cm(11.5)
    col_w = Cm(14)

    # Left column: DOPE Hubs
    text_box(s, Cm(2), bottom_y, col_w, Cm(0.9),
             "DOPE Hubs (Non-Profit)",
             font_size=20, bold=True, color=LIGHT_BLUE, font_name=FONT_DISPLAY)
    text_box(s, Cm(2), bottom_y + Cm(1.2), col_w, Cm(2.5),
             "Excellence training and applied research.\nThe product is the people.",
             font_size=13, color=WHITE, font_name=FONT_MONO, line_spacing=1.5)

    # Arrow between columns (positioned above text to avoid overlap)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(16.5), bottom_y + Cm(0.3), Cm(1.0), Cm(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()

    # Right column: Firmamento Technologies
    text_box(s, Cm(18), bottom_y, col_w, Cm(0.9),
             "Firmamento Technologies",
             font_size=20, bold=True, color=GOLD, font_name=FONT_DISPLAY)
    text_box(s, Cm(18), bottom_y + Cm(0.9), col_w, Cm(0.5),
             "Deep-Tech Cooperative",
             font_size=11, italic=True, color=GOLD_DARK, font_name=FONT_MONO)
    text_box(s, Cm(18), bottom_y + Cm(1.6), col_w, Cm(2.5),
             "Industrialization and Venture Building.\nThe product is digital sovereignty.",
             font_size=13, color=WHITE, font_name=FONT_MONO, line_spacing=1.5)

    # KPI strip
    kpi_y = Cm(15.5)
    add_line(s, Cm(1.5), kpi_y, SW - Cm(1.5), kpi_y, color=GOLD, weight=0.6)

    kpi_data = [
        ("20 → 120", "active researchers in 10 months"),
        ("9", "parallel projects, Aerospace, Health, AI, Cyber, EO"),
        ("6th worldwide", "UAS Challenge 2025 + Best Newcomer Award"),
        ("Institutional", "University of Genoa, Liguria Region, RINA, Legacoop"),
    ]
    kpi_x = Cm(1.5)
    kpi_w = (SW - Cm(3)) / 4
    for value, desc in kpi_data:
        text_box(s, kpi_x, kpi_y + Cm(0.3), kpi_w - Cm(0.3), Cm(0.9),
                 value,
                 font_size=18, bold=True, color=GOLD, font_name=FONT_DISPLAY,
                 align=PP_ALIGN.LEFT)
        text_box(s, kpi_x, kpi_y + Cm(1.4), kpi_w - Cm(0.3), Cm(1.5),
                 desc,
                 font_size=9, color=WHITE, font_name=FONT_MONO,
                 align=PP_ALIGN.LEFT, line_spacing=1.4)
        kpi_x += kpi_w

    # Footer
    text_box(s, Cm(1.5), SH - Cm(0.7), SW - Cm(3), Cm(0.5),
             "FIRMAMENTO TECHNOLOGIES SOCIETA' COOPERATIVA . P.IVA IT03038500991 . PEC firmamentotechnologies@pec.it",
             font_size=7, color=GREY_LIGHT, font_name=FONT_MONO,
             align=PP_ALIGN.CENTER)


# ============== SLIDE 2: PROJECT HALE (intro + use cases) ==============
def slide_2_hale_intro():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, fill=NAVY)

    # Title
    text_box(s, Cm(1.5), Cm(1.2), Cm(20), Cm(1.8),
             "Project H.A.L.E.",
             font_size=40, bold=True, color=WHITE, font_name=FONT_DISPLAY)
    text_box(s, Cm(1.5), Cm(3.0), Cm(20), Cm(0.8),
             "(High Altitude Long Endurance)",
             font_size=18, color=WHITE, font_name=FONT_DISPLAY)

    text_box(s, Cm(1.5), Cm(4.6), Cm(15), Cm(2.0),
             "The cooperative pseudo-satellite\nfor national resilience.",
             font_size=16, color=WHITE, font_name=FONT_MONO, line_spacing=1.4)

    # Tech specs
    specs_y = Cm(8)
    specs = [
        ("Operating altitude", "20 km"),
        ("Macro-regional coverage", "> 500 km"),
        ("Latency", "< 20 ms"),
        ("Propulsion", "Solar-electric"),
        ("Endurance target", "30+ days perennial (Y3)"),
    ]
    for label, val in specs:
        rich_text_box(s, Cm(1.5), specs_y, Cm(15), Cm(0.9), [
            {"text": label + " ", "size": 13, "color": LIGHT_BLUE, "font": FONT_MONO},
            {"text": val, "size": 13, "color": WHITE, "bold": True, "font": FONT_MONO},
        ])
        specs_y += Cm(0.85)

    # HALE Render right (medium size)
    add_image(s, HALE_RENDER, Cm(16), Cm(3.5), w=Cm(12.5), h=Cm(8.5))

    # Use cases right column
    uc_y = Cm(13)
    uc_x = Cm(1.5)
    uc_w = SW - Cm(3)
    text_box(s, uc_x, uc_y, uc_w, Cm(0.8),
             "Public value, four mission profiles",
             font_size=15, bold=True, color=GOLD, font_name=FONT_DISPLAY)

    cases = [
        ("Disaster Response", "Rapid network restoration after floods and landslides."),
        ("Environmental Monitoring", "Wildfire alert and hydrogeological surveillance (multispectral / thermal)."),
        ("Telemedicine", "Ultra-reliable links for remote health posts."),
        ("Civic Connectivity", "Smart agriculture and local e-governance support."),
    ]
    case_w = (uc_w - Cm(1.5)) / 4
    case_x = uc_x
    for title, desc in cases:
        text_box(s, case_x, uc_y + Cm(1.1), case_w, Cm(0.7),
                 title,
                 font_size=12, bold=True, color=LIGHT_BLUE, font_name=FONT_DISPLAY)
        text_box(s, case_x, uc_y + Cm(2.0), case_w, Cm(2.5),
                 desc,
                 font_size=10, color=WHITE, font_name=FONT_MONO, line_spacing=1.4)
        case_x += case_w + Cm(0.5)

    # Footer
    add_line(s, Cm(1.5), SH - Cm(1.4), SW - Cm(1.5), SH - Cm(1.4), color=GOLD, weight=0.5)
    text_box(s, Cm(1.5), SH - Cm(1.0), Cm(15), Cm(0.5),
             "FIRMAMENTO TECHNOLOGIES",
             font_size=9, bold=True, color=GOLD, font_name=FONT_DISPLAY)
    text_box(s, SW - Cm(16.5), SH - Cm(1.0), Cm(15), Cm(0.5),
             "Pentema pilot . Liguria SNAI . H.A.L.E. Stratospheric Layer",
             font_size=9, color=WHITE, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)


# ============== SLIDE 3: BEYOND THE DRONE ==============
def slide_3_beyond():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, fill=WHITE)

    # Title
    text_box(s, Cm(1.5), Cm(1.2), Cm(30), Cm(1.8),
             "Beyond the Drone",
             font_size=36, bold=True, color=NAVY, font_name=FONT_DISPLAY)
    text_box(s, Cm(1.5), Cm(3.1), Cm(30), Cm(0.9),
             "Three structural pillars defining the strategic positioning.",
             font_size=15, color=GREY, font_name=FONT_MONO)

    # Three columns
    cols = [
        ("Data\nSovereignty", LIGHT_BLUE,
         "European data residency. Consent management and sovereign cloud integration. Full GDPR + NIS2 compliance, AgID/PSN hosting for PA-grade data."),
        ("Dual-Use\nResilience", GOLD,
         "Support and backup to critical infrastructure. Persistent regional coverage for civil protection, disaster recovery and emergency telecom."),
        ("Seasonal\nPersistence", LIGHT_BLUE,
         "Architecture optimized for the 44 deg N energy balance, with continuous monitoring March-October and seasonal-only fallback as Plan A."),
    ]
    col_x = Cm(1.5)
    col_y = Cm(5.5)
    col_w = Cm(10)
    col_h = Cm(8)
    spacing = Cm(0.5)
    for title, color, desc in cols:
        text_box(s, col_x, col_y, col_w, Cm(3),
                 title,
                 font_size=32, bold=True, color=color, font_name=FONT_DISPLAY,
                 line_spacing=1.05)
        text_box(s, col_x, col_y + Cm(3.5), col_w, col_h - Cm(3.5),
                 desc,
                 font_size=13, color=NAVY, font_name=FONT_MONO, line_spacing=1.4)
        col_x += col_w + spacing

    # Bottom box: pilot + numbers
    bottom_y = Cm(14)
    bottom_h = Cm(3.5)
    add_rect(s, Cm(1.5), bottom_y, SW - Cm(3), bottom_h, fill=NAVY)

    rich_text_box(s, Cm(2), bottom_y + Cm(0.4), Cm(13), Cm(1),
        [
            {"text": "Pilot site . ", "size": 11, "color": GOLD, "font": FONT_MONO, "bold": True},
            {"text": "Pentema, Torriglia (Genova), Liguria, Italy . SNAI Inner Area",
             "size": 11, "color": WHITE, "font": FONT_MONO},
        ])
    rich_text_box(s, Cm(2), bottom_y + Cm(1.0), Cm(13), Cm(1),
        [
            {"text": "Network . ", "size": 11, "color": GOLD, "font": FONT_MONO, "bold": True},
            {"text": "10 Legacoop cooperatives . lead Fabrica",
             "size": 11, "color": WHITE, "font": FONT_MONO},
        ])
    rich_text_box(s, Cm(2), bottom_y + Cm(1.6), Cm(13), Cm(1),
        [
            {"text": "Tender . ", "size": 11, "color": GOLD, "font": FONT_MONO, "bold": True},
            {"text": "Coopfond Cooding Prototypes (Legacoop)",
             "size": 11, "color": WHITE, "font": FONT_MONO},
        ])
    rich_text_box(s, Cm(2), bottom_y + Cm(2.2), Cm(13), Cm(1),
        [
            {"text": "Compliance . ", "size": 11, "color": GOLD, "font": FONT_MONO, "bold": True},
            {"text": "Italian PFTE art. 41 D.Lgs. 36/2023 + NASA SE Handbook Rev 2",
             "size": 11, "color": WHITE, "font": FONT_MONO},
        ])

    nums_x = Cm(17)
    nums_w = Cm(13)
    rich_text_box(s, nums_x, bottom_y + Cm(0.4), nums_w, Cm(1.2),
        [
            {"text": "€700k - €2M ", "size": 18, "color": GOLD, "bold": True, "font": FONT_DISPLAY},
            {"text": "CapEx Y1", "size": 12, "color": WHITE, "font": FONT_MONO},
        ])
    rich_text_box(s, nums_x, bottom_y + Cm(1.4), nums_w, Cm(1.2),
        [
            {"text": "€1.18M/y ", "size": 18, "color": GOLD, "bold": True, "font": FONT_DISPLAY},
            {"text": "OpEx Y2 reconciled (incl. regulatory team)", "size": 12, "color": WHITE, "font": FONT_MONO},
        ])
    rich_text_box(s, nums_x, bottom_y + Cm(2.4), nums_w, Cm(1.2),
        [
            {"text": "€260k ", "size": 18, "color": GOLD, "bold": True, "font": FONT_DISPLAY},
            {"text": "Revenue Y1 baseline (range €220-300k)", "size": 12, "color": WHITE, "font": FONT_MONO},
        ])

    # Footer
    text_box(s, Cm(1.5), SH - Cm(1.0), Cm(15), Cm(0.5),
             "FIRMAMENTO TECHNOLOGIES SOCIETA' COOPERATIVA",
             font_size=8, bold=True, color=NAVY, font_name=FONT_DISPLAY)
    text_box(s, SW - Cm(16.5), SH - Cm(1.0), Cm(15), Cm(0.5),
             "P.IVA IT03038500991 . PEC firmamentotechnologies@pec.it",
             font_size=8, color=GREY, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)


def build():
    slide_1_ecosystem()
    slide_2_hale_intro()
    slide_3_beyond()
    prs.save(OUT_PPTX)
    print(f"PPTX generato: {OUT_PPTX}")
    print(f"Size: {os.path.getsize(OUT_PPTX)/1024:.1f} KB")


if __name__ == "__main__":
    build()
