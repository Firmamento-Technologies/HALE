"""
Build PPTX deck per CdA Firmamento + sponsor istituzionali (Coopfond, Regione Liguria).
~18 slide riassunto Studio HALE/VTOL.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

BASE = "/home/user/HALE/studio-di-fattibilita"
ASSETS = "/home/user/HALE/cad"
OUT = os.path.join(BASE, "deliverables-FINAL")
LOGO = os.path.join(ASSETS, "LogoFirmamento Technologies.png")
HALE_RENDER = os.path.join(ASSETS, "HALE2.png")
OUT_PPTX = os.path.join(OUT, "DECK-CdA-Sponsor.pptx")

NAVY = RGBColor(0x0a, 0x1a, 0x3d)
GOLD = RGBColor(0xf0, 0xc9, 0x5c)
GREY = RGBColor(0x3a, 0x3a, 0x3a)
GREY_LIGHT = RGBColor(0xf4, 0xf4, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)
RED = RGBColor(0xc0, 0x33, 0x33)
GREEN_OK = RGBColor(0x2e, 0x7d, 0x32)
YELLOW_W = RGBColor(0xed, 0x9a, 0x00)

# Setup presentation 16:9 widescreen
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ============== HELPER FUNCTIONS ==============
def add_text_box(slide, x, y, w, h, text, font_size=18, bold=False, italic=False,
                 color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font_name
    return tb

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)

def add_master_header(slide, title, slide_num, total):
    """Bandiera navy in alto + logo + numero slide."""
    # Banda top navy
    add_rect(slide, 0, 0, SW, Cm(1.2), fill=NAVY)
    # Banda gold sotto
    add_rect(slide, 0, Cm(1.2), SW, Cm(0.12), fill=GOLD)
    # Logo header
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Cm(0.4), Cm(0.15),
                                  height=Cm(0.95))
    # Titolo slide
    add_text_box(slide, Cm(6.5), Cm(0.25), Cm(20), Cm(0.85),
                 title, font_size=18, bold=True,
                 color=WHITE, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Numero slide
    add_text_box(slide, Cm(29), Cm(0.25), Cm(4.5), Cm(0.85),
                 f"{slide_num} / {total}", font_size=11,
                 color=GOLD, italic=True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def add_master_footer(slide):
    """Footer con dati Firmamento."""
    add_rect(slide, 0, SH - Cm(0.7), SW, Cm(0.05), fill=GOLD)
    add_text_box(slide, Cm(0.5), SH - Cm(0.6), SW - Cm(1), Cm(0.5),
                 "Firmamento Technologies Soc. Coop. . P.IVA IT03038500991 . Via Brigata Liguria 105 R, Genova . PEC firmamentotechnologies@pec.it",
                 font_size=8, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 1: COVER ==============
def slide_cover(idx, total):
    s = prs.slides.add_slide(BLANK)
    # Background navy
    add_rect(s, 0, 0, SW, SH, fill=NAVY)
    # Accent gold strip
    add_rect(s, 0, Cm(0), Cm(0.3), SH, fill=GOLD)
    add_rect(s, SW - Cm(0.3), 0, Cm(0.3), SH, fill=GOLD)
    # Logo grande centrato
    add_image(s, LOGO, Cm(11.5), Cm(2), w=Cm(11), h=Cm(5))
    # Titolo
    add_text_box(s, Cm(2), Cm(8), SW - Cm(4), Cm(2.5),
                 "Studio di Fattibilita Tecnico-Economica",
                 font_size=36, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), Cm(10.5), SW - Cm(4), Cm(2),
                 "Piattaforma Aerea HALE / VTOL per le Aree Interne Italiane",
                 font_size=22, italic=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Render HALE
    add_image(s, HALE_RENDER, Cm(10), Cm(12.5), w=Cm(14), h=Cm(4.5))
    # Footer cover
    add_text_box(s, Cm(2), SH - Cm(1.5), SW - Cm(4), Cm(0.6),
                 "FIRMAMENTO TECHNOLOGIES SOCIETA' COOPERATIVA . Bozza M+11 . Maggio 2026",
                 font_size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), SH - Cm(0.9), SW - Cm(4), Cm(0.4),
                 "Deck per Consiglio di Amministrazione + Sponsor istituzionali (Coopfond, Regione Liguria)",
                 font_size=9, italic=True, color=GREY_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============== STD SLIDE ==============
def std_slide(title, idx, total):
    s = prs.slides.add_slide(BLANK)
    add_master_header(s, title, idx, total)
    add_master_footer(s)
    return s

# ============== SLIDE 2: AGENDA ==============
def slide_agenda(idx, total):
    s = std_slide("Agenda", idx, total)
    items = [
        ("1.", "Il progetto in pillole"),
        ("2.", "Strategia duale, Percorso 6A e Percorso 6B"),
        ("3.", "Verdetto consolidato"),
        ("4.", "Caso pilota Pentema, contesto territoriale"),
        ("5.", "Modello di business e value proposition"),
        ("6.", "Quadro economico-finanziario"),
        ("7.", "Mix finanziamento raccomandato Y1"),
        ("8.", "Cronoprogramma e gate decisionali"),
        ("9.", "Risk profile e showstopper"),
        ("10.", "Visione 10 anni e posizionamento sovrano EU"),
        ("11.", "Decisione formale richiesta al CdA"),
        ("12.", "Action items Y0-Y1 prioritari"),
    ]
    y = Cm(2.5)
    for num, txt in items:
        add_text_box(s, Cm(2.5), y, Cm(2), Cm(0.9), num,
                     font_size=18, bold=True, color=GOLD,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(5), y, SW - Cm(7), Cm(0.9), txt,
                     font_size=16, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.05)

# ============== SLIDE 3: IL PROGETTO ==============
def slide_progetto(idx, total):
    s = std_slide("Il progetto in 5 punti", idx, total)
    points = [
        ("Cosa",
         "Piattaforma aerea unmanned per servizi territoriali alle Aree Interne italiane, operata come operatore di servizi ricorrenti, non come venditore di velivoli."),
        ("Dove pilota",
         "Pentema, frazione del Comune di Torriglia (Genova), 14 abitanti ISTAT, area SNAI Valli Antola-Tigullio."),
        ("Strategia duale",
         "Percorso 6A VTOL pilota (0-12 mesi, low-risk) + Percorso 6B HALE stratosferico R&D Phase B (24-48+ mesi, high-risk)."),
        ("Use case",
         "Monitoraggio frane, antincendio boschivo, connettivita di emergenza, mapping infrastrutture, agricoltura cooperative."),
        ("Visione 10 anni",
         "Nodo italiano fondatore di una futura infrastruttura sovrana europea HAPS, complementare a IRIS²."),
    ]
    y = Cm(2.8)
    for label, desc in points:
        # Box label gold
        add_rect(s, Cm(2), y, Cm(5), Cm(2), fill=GOLD, line=NAVY)
        add_text_box(s, Cm(2), y, Cm(5), Cm(2), label,
                     font_size=15, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Desc
        add_text_box(s, Cm(7.5), y, SW - Cm(9.5), Cm(2), desc,
                     font_size=13, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(2.4)

# ============== SLIDE 4: STRATEGIA DUALE ==============
def slide_strategia(idx, total):
    s = std_slide("Strategia duale, Percorso 6A + 6B", idx, total)
    # 6A box
    add_rect(s, Cm(1.5), Cm(2.5), Cm(15), Cm(13), fill=GREY_LIGHT, line=NAVY)
    add_rect(s, Cm(1.5), Cm(2.5), Cm(15), Cm(1.2), fill=NAVY)
    add_text_box(s, Cm(1.5), Cm(2.5), Cm(15), Cm(1.2),
                 "Percorso 6A, VTOL Pilota Pentema",
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    items_6a = [
        ("Tecnologia", "VTOL ibrido TRL 8-9 (JOUAV CW-30E + Plan B Tekever AR3)"),
        ("Orizzonte", "0-12 mesi (M+0 → M+12)"),
        ("Budget CapEx Y1", "€700k - €2M (IVA + contingency 15% incluse)"),
        ("OpEx Y2 RECONCILED", "€1.18M/anno centrale (con +3 FTE regulatory team)"),
        ("Revenue Y1 RECALIBRATED", "€260k centrale (range €220-300k, min €200k SyR-Cost-003)"),
        ("Casi d'uso", "Frane, antincendio, connettivita emergenza, mapping, cooperative"),
        ("Verdetto", "HOLD CON PIANO REGOLATORIO RAFFORZATO (P 45-60%) / GO COND. (P 5-15%)"),
    ]
    y = Cm(4.2)
    for label, val in items_6a:
        add_text_box(s, Cm(2), y, Cm(5), Cm(0.7), label,
                     font_size=10, bold=True, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(7), y, Cm(9), Cm(0.7), val,
                     font_size=10, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.25)

    # 6B box
    add_rect(s, Cm(17.5), Cm(2.5), Cm(15), Cm(13), fill=GREY_LIGHT, line=NAVY)
    add_rect(s, Cm(17.5), Cm(2.5), Cm(15), Cm(1.2), fill=NAVY)
    add_text_box(s, Cm(17.5), Cm(2.5), Cm(15), Cm(1.2),
                 "Percorso 6B, HALE Stratosferico R&D",
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    items_6b = [
        ("Tecnologia", "HALE solare, apertura 25-30 m, MTOW 80-150 kg"),
        ("Quota operativa", "18-21 km (FL590-690)"),
        ("Orizzonte", "24-48+ mesi (Phase B R&D)"),
        ("Budget Phase B", "€5.5-13.5M (mix 50-75% grant pubblico)"),
        ("Pivot strutturale M+3", "Operatore di servizi su piattaforme prime contractor"),
        ("Showstopper aperti", "5 RSK rossi (energy, aeroelast., HAPS framework, funding, TC)"),
        ("Verdetto", "HOLD CON CRITERI USCITA STRINGENTI + Pivot Strutturale"),
    ]
    y = Cm(4.2)
    for label, val in items_6b:
        add_text_box(s, Cm(18), y, Cm(5), Cm(0.7), label,
                     font_size=10, bold=True, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(23), y, Cm(9), Cm(0.7), val,
                     font_size=10, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.25)

# ============== SLIDE 5: VERDETTO ==============
def slide_verdetto(idx, total):
    s = std_slide("Verdetto consolidato", idx, total)
    # Table verdetto
    data = [
        ["Percorso", "Verdetto", "Probabilita", "Razionale"],
        ["6A VTOL pilota",
         "HOLD CON PIANO RAFFORZATO (base)\nGO CONDIZIONATO (ottimistico)",
         "45-60% Hold\n5-15% Go pieno",
         "Solido tecnicamente, debole su engagement esterno"],
        ["6B HALE R&D",
         "HOLD CON CRITERI USCITA STRINGENTI\n+ Pivot strutturale prime contractor",
         "Hold permanente\nfino gate G5 M+24",
         "Base rate 0% HALE solari commerciali in 22 anni; capital intensity $50M-1B benchmark"],
    ]
    # Header navy
    y0 = Cm(2.6)
    col_x = [Cm(1.5), Cm(7), Cm(16), Cm(22)]
    col_w = [Cm(5.5), Cm(9), Cm(6), Cm(10.4)]
    row_h = [Cm(0.9), Cm(3.5), Cm(3.5)]

    for ci, (x, w, txt) in enumerate(zip(col_x, col_w, data[0])):
        add_rect(s, x, y0, w, row_h[0], fill=NAVY)
        add_text_box(s, x, y0, w, row_h[0], txt,
                     font_size=13, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    y = y0 + row_h[0]
    for ri, row in enumerate(data[1:]):
        row_color = WHITE if ri % 2 == 0 else GREY_LIGHT
        for ci, (x, w, txt) in enumerate(zip(col_x, col_w, row)):
            add_rect(s, x, y, w, row_h[ri + 1], fill=row_color, line=NAVY)
            color = NAVY if ci == 0 else (RED if ri == 1 and ci == 1 else (YELLOW_W if ci == 1 else GREY))
            add_text_box(s, x, y, w, row_h[ri + 1], txt,
                         font_size=11, bold=(ci == 0 or ci == 1), color=color,
                         align=PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
        y += row_h[ri + 1]

    # Caveat strip
    add_rect(s, Cm(1.5), Cm(14.5), Cm(30.9), Cm(2.5), fill=RGBColor(0xfd, 0xf8, 0xe8), line=GOLD)
    add_text_box(s, Cm(2), Cm(14.6), Cm(30), Cm(2.3),
                 "CAVEAT PROBABILISTICO ONESTO (post audit M+3): le 5 hard conditions C1-C5 del Percorso 6A sono in AND logico. "
                 "P(AND tutte soddisfatte) realistica al gate G3 e 5-15% per Go pieno, 45-60% per Hold con piano e re-review M+13-16. "
                 "Scenario base atteso = HOLD CON PIANO RAFFORZATO.",
                 font_size=11, italic=True, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 6: PENTEMA ==============
def slide_pentema(idx, total):
    s = std_slide("Caso pilota Pentema, contesto territoriale", idx, total)
    # Box dati
    data = [
        ("Localita", "Pentema, frazione del Comune di Torriglia (GE)"),
        ("Coordinate", "Circa 44.5°N 9.2°E"),
        ("Altitudine", "1100-1300 m s.l.m."),
        ("Popolazione ISTAT", "14 residenti (7 M, 7 F), 11 famiglie, 100 edifici"),
        ("Distanza centro abitato Torriglia", "3.27 km"),
        ("Area SNAI", "Valli Antola-Tigullio (ciclo 2021-2027)"),
        ("Vincolo ambientale", "Parco Naturale Regionale dell'Antola (L.R. Liguria 12/1995)"),
        ("Rete Natura 2000", "SIC/ZSC/ZPS IT1331402"),
        ("Specie target Allegato I", "Aquila chrysaetos, Bubo bubo (avifauna nidificante)"),
        ("Rete utenti pilota", "10 cooperative Legacoop, capofila Fabrica"),
    ]
    y = Cm(2.5)
    for label, val in data:
        add_text_box(s, Cm(2), y, Cm(11), Cm(0.85), label,
                     font_size=12, bold=True, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(13.5), y, Cm(18), Cm(0.85), val,
                     font_size=12, color=GREY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.1)

# ============== SLIDE 7: BUSINESS MODEL ==============
def slide_business(idx, total):
    s = std_slide("Modello di business, service-only su cooperative", idx, total)
    # 4 quadranti
    quadrants = [
        (Cm(1.5), Cm(2.5), Cm(15), Cm(6.5),
         "Service-only", "Nessuna vendita di velivoli. Erogazione di servizi ricorrenti: monitoraggio EO, connettivita emergenza, analytics, ore-volo + analytics, outcome-based, DaaS."),
        (Cm(17.5), Cm(2.5), Cm(15), Cm(6.5),
         "4 pilastri vantaggio competitivo",
         "(i) specializzazione geografica Aree Interne, (ii) modello cooperativo Legacoop, (iii) sostenibilita ESG (propulsione solare/elettrica + fibra lino), (iv) approccio incrementale VTOL → HALE."),
        (Cm(1.5), Cm(9.5), Cm(15), Cm(6.5),
         "Canali distributivi",
         "B2G regionale (anchor Regione Liguria) - B2G locale (Protezione Civile, ARPA, Enti Parco) - B2B cooperative (rete Legacoop scaled)."),
        (Cm(17.5), Cm(9.5), Cm(15), Cm(6.5),
         "Pricing PA RECALIBRATED",
         "Post audit Cluster D (Planetek, e-GEOS, NHazca): €60-90k/anno base + €30-60k premium persistence/sovranita. Falsificato il baseline originale €150k/anno → €355-405k revenue Y1."),
    ]
    for x, y, w, h, title, txt in quadrants:
        add_rect(s, x, y, w, h, fill=GREY_LIGHT, line=NAVY)
        add_rect(s, x, y, w, Cm(1.2), fill=NAVY)
        add_text_box(s, x, y, w, Cm(1.2), title,
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, x + Cm(0.4), y + Cm(1.4), w - Cm(0.8), h - Cm(1.6), txt,
                     font_size=12, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# ============== SLIDE 8: NUMERI CHIAVE ==============
def slide_numeri(idx, total):
    s = std_slide("Numeri chiave del progetto", idx, total)
    data = [
        ["Metrica", "Valore", "Note"],
        ["CapEx Y1 (6A)", "€700k - €2M", "IVA + contingency 15% incluse; sliding realistic €2.5-3.5M"],
        ["OpEx Y2 RECONCILED", "€1.18M/anno (centro)", "Range €1.05-1.30M con +3 FTE regulatory mandatory"],
        ["Revenue Y1 RECALIBRATED", "€260k (centro)", "Range €220-300k, hard floor €200k SyR-Cost-003"],
        ["Break-even cumulato", "Y5-Y6", "Post recalibration revenue + OpEx"],
        ["NPV 10y scenario base", "+€3.5M", "WACC 12% blended (grant + equity + R&D credit)"],
        ["IRR 10y scenario base", "12-18%", "Post recalibration; pre era 18-25%"],
        ["Payback semplice", "5-7 anni", "Post recalibration"],
        ["ARR Y3 target", "€1.5-3.5M", "Scale-up Liguria + 1 regione"],
        ["ARR Y5 target", "€3-8M", "Multi-regione + HAPS subscale"],
        ["Capital intensity Y10 small fleet", "€500M - €2B", "5-10 HAPS + 10-20 VTOL/MALE"],
        ["Capital intensity Y10 EU sovereign", "€10-30B", "Precondizione: programma EU dedicato analog IRIS²"],
    ]
    y = Cm(2.5)
    col_x = [Cm(1.5), Cm(15), Cm(20)]
    col_w = [Cm(13.5), Cm(5), Cm(12.4)]
    row_h = Cm(1.0)
    for ri, row in enumerate(data):
        for ci, (x, w, txt) in enumerate(zip(col_x, col_w, row)):
            if ri == 0:
                add_rect(s, x, y, w, row_h, fill=NAVY)
                add_text_box(s, x, y, w, row_h, txt,
                             font_size=12, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                row_color = WHITE if ri % 2 == 1 else GREY_LIGHT
                add_rect(s, x, y, w, row_h, fill=row_color, line=GREY)
                add_text_box(s, x + Cm(0.2), y, w - Cm(0.4), row_h, txt,
                             font_size=10, bold=(ci == 0 or ci == 1),
                             color=NAVY if ci < 2 else GREY,
                             align=PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.LEFT,
                             anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

# ============== SLIDE 9: MIX FUNDING ==============
def slide_funding(idx, total):
    s = std_slide("Mix finanziamento raccomandato Y1 Percorso 6A", idx, total)
    data = [
        ["Fonte", "Target €", "% mix", "Status M+3"],
        ["Coopfond Cooding Prototypes 2026", "50k", "5%", "DR-002 verifica calendario bando"],
        ["Coopfond Cooding-Invest", "150-300k", "15-20%", "Q2 2026 pending"],
        ["Regione Liguria FESR 2021-2027", "300-500k", "25-40%", "OQ-010 LoI in costruzione"],
        ["PNRR Aerospazio / IS4Aerospace", "0-300k", "0-20%", "Partnership Polito DIMEAS in valutazione"],
        ["Equity privato (founder + seed)", "200-500k", "15-35%", "Round seed Q1 2026"],
        ["R&D tax credit (L. 160/2019)", "50-150k", "5-15%", "Cumulabile post-spesa"],
        ["TOTALE", "€0.75-1.75M", "100%", "Mix da consolidare entro M+10 (60% committed)"],
    ]
    y = Cm(2.5)
    col_x = [Cm(1.5), Cm(15), Cm(19), Cm(23)]
    col_w = [Cm(13.5), Cm(4), Cm(4), Cm(9.4)]
    row_h = Cm(1.4)
    for ri, row in enumerate(data):
        for ci, (x, w, txt) in enumerate(zip(col_x, col_w, row)):
            is_total = (ri == len(data) - 1)
            if ri == 0:
                add_rect(s, x, y, w, row_h, fill=NAVY)
                add_text_box(s, x, y, w, row_h, txt,
                             font_size=13, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            elif is_total:
                add_rect(s, x, y, w, row_h, fill=GOLD)
                add_text_box(s, x + Cm(0.2), y, w - Cm(0.4), row_h, txt,
                             font_size=12, bold=True, color=NAVY,
                             align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
                             anchor=MSO_ANCHOR.MIDDLE)
            else:
                row_color = WHITE if ri % 2 == 1 else GREY_LIGHT
                add_rect(s, x, y, w, row_h, fill=row_color, line=GREY)
                add_text_box(s, x + Cm(0.2), y, w - Cm(0.4), row_h, txt,
                             font_size=11, color=NAVY if ci < 3 else GREY,
                             bold=(ci == 0),
                             align=PP_ALIGN.CENTER if ci in [1, 2] else PP_ALIGN.LEFT,
                             anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

# ============== SLIDE 10: CRONOPROGRAMMA ==============
def slide_crono(idx, total):
    s = std_slide("Cronoprogramma e gate decisionali", idx, total)
    # Gate timeline visiva
    gates = [
        ("G0", "M+0", "Kick-off", NAVY),
        ("G1", "M+3", "Concept Review", NAVY),
        ("G2", "M+6", "Architecture Baselined", NAVY),
        ("G3", "M+10/11", "★ FEASIBILITY GATE", GOLD),
        ("G4", "M+12", "Fine pilota VTOL", NAVY),
        ("G5", "M+24", "Phase B 6B Decision", NAVY),
        ("G6", "M+36", "Phase B Midterm", NAVY),
        ("End", "M+48", "Phase B End", NAVY),
    ]
    # Linea timeline
    add_rect(s, Cm(1.5), Cm(8), Cm(30.9), Cm(0.15), fill=NAVY)
    # Gate boxes
    box_w = Cm(3.7)
    spacing = Cm(0.2)
    total_w = len(gates) * box_w + (len(gates) - 1) * spacing
    start_x = (SW - total_w) / 2

    for i, (gid, mese, label, color) in enumerate(gates):
        x = start_x + i * (box_w + spacing)
        is_g3 = (gid == "G3")
        # Vertical line
        add_rect(s, x + box_w / 2 - Cm(0.05), Cm(7), Cm(0.1), Cm(2), fill=color)
        # Box gate
        h_box = Cm(3) if is_g3 else Cm(2.5)
        y_box = Cm(9) if is_g3 else Cm(9.5)
        add_rect(s, x, y_box, box_w, h_box, fill=color)
        add_text_box(s, x, y_box + Cm(0.3), box_w, Cm(0.8), gid,
                     font_size=18 if is_g3 else 15, bold=True,
                     color=NAVY if is_g3 else WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, x, y_box + Cm(1.0), box_w, Cm(0.6), mese,
                     font_size=11, bold=True,
                     color=NAVY if is_g3 else WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, x, y_box + Cm(1.5), box_w, h_box - Cm(1.5), label,
                     font_size=9, bold=is_g3,
                     color=NAVY if is_g3 else WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Fase labels above
    add_text_box(s, Cm(2), Cm(3.5), Cm(15), Cm(1), "STUDIO DI FATTIBILITA' (M+0 → M+11)",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(15.5), Cm(3.5), Cm(8), Cm(1), "PILOTA 6A VTOL (M+12 → M+24)",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(23), Cm(3.5), Cm(9), Cm(1), "R&D 6B HALE PHASE B (M+24 → M+48)",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # G3 highlight
    add_rect(s, Cm(1.5), Cm(14), Cm(30.9), Cm(2.5), fill=RGBColor(0xfd, 0xf8, 0xe8), line=GOLD)
    add_text_box(s, Cm(2), Cm(14.1), Cm(30), Cm(2.3),
                 "★ G3 FEASIBILITY GATE PRIMARIO (M+10-11): verdetto Go / Hold / No-Go per ciascun percorso. "
                 "Oggetto formale del presente Studio. Scenario realistico atteso: HOLD CON PIANO RAFFORZATO + re-review G3-bis M+13-16.",
                 font_size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 11: RISK PROFILE ==============
def slide_risk(idx, total):
    s = std_slide("Risk profile, 5 showstopper formali (Risk Register A.2)", idx, total)
    data = [
        ["#", "RSK-ID", "Categoria", "Score", "Mitigazione"],
        ["1", "RSK-TEC-001", "Tecnico, energy balance HALE inverno 44°N", "25", "E5 Seasonal-only marzo-ottobre mandatory Plan A; -50.1% deficit confermato"],
        ["2", "RSK-REG-001", "Regolatorio, framework HAPS EU mancante", "20", "Engagement EASA Innovation Network + RMT HAPS request"],
        ["3", "RSK-FIN-001", "Finanziario, Phase B €5.5-13.5M non committed", "20", "Mix funding 50%+ pubblico + Series A/B €5-15M raised"],
        ["4", "RSK-TEC-002", "Tecnico, aeroelasticita ala high-AR", "15", "Aeroelastic analysis preliminare M+10"],
        ["5", "RSK-TEC-003", "Tecnico, Type Certification HALE > 5 anni", "16", "Special Condition negoziata caso per caso EASA"],
    ]
    y = Cm(2.6)
    col_x = [Cm(1.5), Cm(2.5), Cm(5.5), Cm(15.5), Cm(18.5)]
    col_w = [Cm(1), Cm(3), Cm(10), Cm(3), Cm(13.9)]
    row_h = Cm(1.4)
    for ri, row in enumerate(data):
        for ci, (x, w, txt) in enumerate(zip(col_x, col_w, row)):
            if ri == 0:
                add_rect(s, x, y, w, row_h, fill=NAVY)
                add_text_box(s, x, y, w, row_h, txt,
                             font_size=11, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                row_color = WHITE if ri % 2 == 1 else GREY_LIGHT
                add_rect(s, x, y, w, row_h, fill=row_color, line=GREY)
                color = NAVY if ci != 3 else RED
                add_text_box(s, x + Cm(0.1), y, w - Cm(0.2), row_h, txt,
                             font_size=10, bold=(ci in [0, 1, 3]),
                             color=color,
                             align=PP_ALIGN.CENTER if ci in [0, 3] else PP_ALIGN.LEFT,
                             anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

    # Risk reduction summary
    add_text_box(s, Cm(1.5), Cm(11.5), Cm(15), Cm(0.7),
                 "Risk reduction baseline → post-mitigation:",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(1.5), Cm(12.3), Cm(30), Cm(2.5),
                 "Pre-mitigation: RED 17 . YELLOW 66 . GREEN 33 (su 116 rischi baseline v1.0). "
                 "Post-mitigation residual: RED 2 . YELLOW 19 . GREEN 95. "
                 "RED reduction 88% (17 → 2). YELLOW reduction 71% (66 → 19). "
                 "Aggiornamento v1.5: integrazione 3 RSK-AMB ambientali post A.12 VIA v2.0 → totale 119 rischi in 14 categorie.",
                 font_size=11, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# ============== SLIDE 12: VISIONE 10 ANNI ==============
def slide_visione(idx, total):
    s = std_slide("Visione 10 anni, posizionamento sovrano EU", idx, total)
    fasi = [
        ("Y1", "Pilota Pentema VTOL", "Pilot validato + revenue €260k", NAVY),
        ("Y2-Y3", "Scale-up SNAI Italia", "3-4 regioni, flotta 3-8 VTOL/MALE + R&D HALE subscale", NAVY),
        ("Y3-Y6", "Primo HALE prototipo italiano", "Servizio commerciale HAPS pilota + Series A-B raised", NAVY),
        ("Y6-Y8", "Costellazione italiana", "3-10 HAPS operativi + servizi NTN + EO persistente", NAVY),
        ("Y8-Y10", "Consorzio EU stratospheric layer", "Italia + FR/DE/ES + posizionamento EU sovereign complementare IRIS²", GOLD),
    ]
    y = Cm(2.8)
    for fase, titolo, desc, color in fasi:
        is_y10 = (fase == "Y8-Y10")
        # Box fase
        add_rect(s, Cm(2), y, Cm(4), Cm(2), fill=color)
        add_text_box(s, Cm(2), y, Cm(4), Cm(2), fase,
                     font_size=18, bold=True,
                     color=NAVY if is_y10 else WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Titolo
        add_text_box(s, Cm(6.5), y, Cm(10), Cm(2), titolo,
                     font_size=14, bold=True, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # Desc
        add_text_box(s, Cm(17), y, Cm(15), Cm(2), desc,
                     font_size=11, color=GREY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(2.3)

    # IRIS² nota
    add_rect(s, Cm(1.5), Cm(14.7), Cm(30.9), Cm(1.8), fill=RGBColor(0xfd, 0xf8, 0xe8), line=GOLD)
    add_text_box(s, Cm(2), Cm(14.8), Cm(30), Cm(1.6),
                 "Linguaggio pubblico raccomandato: \"complementare a IRIS²\", mai \"alternativa europea a Starlink\". "
                 "IRIS² baseline LEO+MEO puro (DR-009 closure), governance SpaceRISE, lancio 2029, ops 2031. "
                 "Finestra strategica Firmamento 2027-2030 per posizionamento stratospheric gap-filler via DG CNECT engagement.",
                 font_size=10, italic=True, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 13: DECISIONE CDA ==============
def slide_decisione(idx, total):
    s = std_slide("Decisione formale richiesta al CdA + sponsor", idx, total)
    decisioni = [
        "Approvare Studio di Fattibilita Vol. 1 + 2 + 3 (versione M+11 bozza)",
        "Approvare GO CONDIZIONATO Percorso 6A con 5 hard conditions M+10",
        "Approvare HOLD Percorso 6B + pivot strutturale verso prime contractor",
        "Approvare budget Y1 €2.5-3.5M (CapEx + OpEx + bridge financing + regulatory team)",
        "Approvare doppio binario di pianificazione (nominale + sliding §9.12)",
        "Approvare assunzione 3 FTE regulatory (CISO + DPO + Head Regulatory) entro M+6",
        "Approvare engagement Regione Liguria + Coopfond + Comune Torriglia + cooperative entro M+6",
        "Approvare re-baseline Gate G3-bis a M+13-16 come opzione legittima",
    ]
    y = Cm(2.6)
    for i, dec in enumerate(decisioni, 1):
        # Checkbox
        add_rect(s, Cm(1.8), y + Cm(0.15), Cm(0.7), Cm(0.7), fill=WHITE, line=NAVY)
        add_text_box(s, Cm(3), y, Cm(1.5), Cm(1), f"{i}.",
                     font_size=14, bold=True, color=GOLD,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(4.5), y, Cm(28), Cm(1), dec,
                     font_size=12, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.15)

    # Firma
    add_text_box(s, Cm(1.5), Cm(13.5), Cm(15), Cm(0.6),
                 "Firma CEO Firmamento: ___________________________",
                 font_size=11, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(17), Cm(13.5), Cm(15), Cm(0.6),
                 "Firma CdA: ___________________________",
                 font_size=11, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(1.5), Cm(14.3), Cm(15), Cm(0.6),
                 "Firma Sponsor Coopfond: ___________________________",
                 font_size=11, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(17), Cm(14.3), Cm(15), Cm(0.6),
                 "Firma Sponsor Regione Liguria: ___________________________",
                 font_size=11, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(1.5), Cm(15.5), Cm(31), Cm(0.7),
                 "Data delibera: __________________ . Documento di approvazione M+11 . Maggio 2026",
                 font_size=11, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============== SLIDE 14: TOP-10 ACTION ==============
def slide_actions(idx, total):
    s = std_slide("Top-10 Action Items prioritari Y0-Y1", idx, total)
    actions = [
        ("1.", "Capital structure resistente (founder ≥ 51% + golden share + anti-acquisition)", "M+6"),
        ("2.", "LoI Regione Liguria firmata con DGR specifica", "M+6"),
        ("3.", "Pre-application meeting ENAC + feedback documentato", "M+3-6"),
        ("4.", "Workshop pubblico comunita Pentema (14 abitanti) + DPIA pubblica", "M+3-6"),
        ("5.", "Assunzione 3 FTE regulatory (CISO + DPO + Head Regulatory)", "M+6-9"),
        ("6.", "Benchmark pricing PA reale (e-GEOS, Planetek, NHazca)", "M+6"),
        ("7.", "Quotation JOUAV + Tekever parallele per CapEx accurate", "M+3"),
        ("8.", "Engagement EASA Innovation Network + RMT HAPS request", "M+9-12"),
        ("9.", "Modello finanziario Excel + DCF + sensitivity + Monte Carlo", "M+6-9"),
        ("10.", "Position paper \"Italian Stratospheric Sovereignty\" + dialogo DG CNECT/DEFIS", "M+9-12"),
    ]
    y = Cm(2.5)
    for num, act, when in actions:
        add_text_box(s, Cm(1.5), y, Cm(1.5), Cm(0.9), num,
                     font_size=14, bold=True, color=GOLD,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Cm(3.5), y, Cm(24), Cm(0.9), act,
                     font_size=12, color=NAVY,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Cm(28), y + Cm(0.05), Cm(4), Cm(0.8), fill=NAVY)
        add_text_box(s, Cm(28), y + Cm(0.05), Cm(4), Cm(0.8), when,
                     font_size=11, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(1.05)

# ============== SLIDE FINALE: GRAZIE ==============
def slide_final(idx, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, fill=NAVY)
    add_rect(s, 0, 0, Cm(0.3), SH, fill=GOLD)
    add_rect(s, SW - Cm(0.3), 0, Cm(0.3), SH, fill=GOLD)
    add_image(s, LOGO, Cm(11), Cm(2), w=Cm(12))
    add_text_box(s, Cm(2), Cm(8.5), SW - Cm(4), Cm(2),
                 "Grazie",
                 font_size=44, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), Cm(11), SW - Cm(4), Cm(1.2),
                 "Per domande e approfondimenti:",
                 font_size=15, italic=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), Cm(12.5), SW - Cm(4), Cm(0.9),
                 "PEC: firmamentotechnologies@pec.it",
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), Cm(13.5), SW - Cm(4), Cm(0.7),
                 "FIRMAMENTO TECHNOLOGIES SOCIETA' COOPERATIVA",
                 font_size=12, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), Cm(14.3), SW - Cm(4), Cm(0.6),
                 "Via Brigata Liguria 105 R . 16121 Genova (GE) . P.IVA IT03038500991 . REA 528629",
                 font_size=10, color=GREY_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Cm(2), SH - Cm(1), SW - Cm(4), Cm(0.5),
                 "Studio di Fattibilita HALE/VTOL . Bozza M+11 . Maggio 2026 . Documento riservato",
                 font_size=9, italic=True, color=GREY_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============== BUILD ==============
def build():
    TOTAL = 14
    slide_cover(1, TOTAL)
    slide_agenda(2, TOTAL)
    slide_progetto(3, TOTAL)
    slide_strategia(4, TOTAL)
    slide_verdetto(5, TOTAL)
    slide_pentema(6, TOTAL)
    slide_business(7, TOTAL)
    slide_numeri(8, TOTAL)
    slide_funding(9, TOTAL)
    slide_crono(10, TOTAL)
    slide_risk(11, TOTAL)
    slide_visione(12, TOTAL)
    slide_decisione(13, TOTAL)
    slide_actions(14, TOTAL)
    slide_final(15, TOTAL)
    prs.save(OUT_PPTX)
    size_mb = os.path.getsize(OUT_PPTX) / 1024 / 1024
    print(f"PPTX generato: {OUT_PPTX} ({size_mb:.2f} MB, 15 slide)")

if __name__ == "__main__":
    build()
